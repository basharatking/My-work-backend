# RunDocs Backend v2.0 — main.py
# KEY FIXES:
# 1. PDF-to-Excel: Tables grouped by HEADER SIGNATURE across ALL pages (NOT page-wise)
# 2. Watermark: Fixed fill_opacity parameter & proper overlay rendering

import io, os, re, zipfile, json, time
import subprocess, tempfile, shutil
from pathlib import Path
from typing import List

from fastapi import FastAPI, File, Form, UploadFile, HTTPException
from fastapi.responses import Response
from fastapi.middleware.cors import CORSMiddleware

try:
    import fitz
    _FITZ_OK = True
except ImportError:
    _FITZ_OK = False

try:
    import pdfplumber
    _PLUMBER_OK = True
except ImportError:
    _PLUMBER_OK = False

try:
    from docx import Document
    from docx.shared import Pt
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    _XL_OK = True
except ImportError:
    _XL_OK = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PInches
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    import requests as _req
    _HF_TOKEN = os.environ.get("HF_TOKEN", "")
    _HF_MODEL = os.environ.get("HF_MODEL", "mistralai/Mistral-7B-Instruct-v0.3")
    _HF_URL   = f"https://api-inference.huggingface.co/models/{_HF_MODEL}"
    _AI_OK    = bool(_HF_TOKEN)
except Exception:
    _AI_OK = False

app = FastAPI(title="RunDocs API", version="2.0.0", docs_url=None, redoc_url=None)
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

MAX_BYTES = 25 * 1024 * 1024


def stream_file(data: bytes, media_type: str, filename: str, extra: dict = None) -> Response:
    headers = {
        "Content-Disposition": f'attachment; filename="{filename}"',
        "Content-Length": str(len(data)),
        "X-Content-Type-Options": "nosniff",
        "Cache-Control": "no-store",
        "Access-Control-Expose-Headers": "X-Original-Size,X-Compressed-Size,X-Savings-Pct",
    }
    if extra:
        headers.update(extra)
    return Response(content=data, media_type=media_type, headers=headers)


async def read_file(upload: UploadFile) -> bytes:
    data = await upload.read()
    if len(data) > MAX_BYTES:
        raise HTTPException(413, "File too large. Free limit is 25 MB.")
    return data


def stem(f) -> str:
    return Path(f or "file").stem


def open_fitz(data: bytes):
    if not _FITZ_OK:
        raise HTTPException(500, "PDF engine not available.")
    return fitz.open(stream=data, filetype="pdf")


def extract_text_fitz(data: bytes, max_pages: int = 40) -> str:
    doc = open_fitz(data)
    pages = min(len(doc), max_pages)
    return "\n\n".join(doc[i].get_text() for i in range(pages))


def ai_call(system: str, prompt: str, max_tokens: int = 1500) -> str:
    if not _AI_OK:
        raise HTTPException(503, "AI Tools unavailable. Add HF_TOKEN in Replit Secrets.")
    full_prompt = f"<s>[INST] {system}\n\n{prompt} [/INST]"
    headers = {"Authorization": f"Bearer {_HF_TOKEN}", "Content-Type": "application/json"}
    payload = {
        "inputs": full_prompt,
        "parameters": {"max_new_tokens": min(max_tokens, 1200), "temperature": 0.4, "top_p": 0.9, "do_sample": True, "return_full_text": False},
        "options": {"wait_for_model": True, "use_cache": False},
    }
    try:
        resp = _req.post(_HF_URL, headers=headers, json=payload, timeout=90)
        if resp.status_code == 503:
            time.sleep(15)
            resp = _req.post(_HF_URL, headers=headers, json=payload, timeout=90)
        if resp.status_code == 401:
            raise HTTPException(503, "Invalid HuggingFace token.")
        if not resp.ok:
            err = resp.json() if resp.content else {}
            raise HTTPException(502, f"AI error: {err.get('error', resp.status_code)}")
        result = resp.json()
        text = result[0].get("generated_text", "") if isinstance(result, list) else result.get("generated_text", "")
        for m in ["[/INST]", "[INST]", "<s>", "</s>"]:
            text = text.replace(m, "")
        return text.strip() or "No result generated."
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(502, f"AI error: {str(e)}")


def _to_roman(n: int) -> str:
    val = [1000,900,500,400,100,90,50,40,10,9,5,4,1]
    sym = ["M","CM","D","CD","C","XC","L","XL","X","IX","V","IV","I"]
    result = ""
    for v, s in zip(val, sym):
        while n >= v:
            result += s; n -= v
    return result.lower()


@app.get("/health")
def health():
    return {"status": "ok", "brand": "RunDocs", "version": "2.0.0",
            "fitz": _FITZ_OK, "plumber": _PLUMBER_OK, "docx": _DOCX_OK,
            "xlsx": _XL_OK, "pptx": _PPTX_OK, "ai": _AI_OK}


# ══════════════════════════════════════════
# AI ENDPOINTS
# ══════════════════════════════════════════

@app.post("/ai-summary")
async def ai_summary(file: UploadFile = File(...), length: str = Form("medium")):
    data = await read_file(file)
    text = extract_text_fitz(data, 20)
    if not text.strip(): raise HTTPException(400, "No readable text found.")
    lens = {"short": "2–3 sentences", "medium": "1 clear paragraph", "long": "2–3 detailed paragraphs"}
    result = ai_call("You are a professional document summarizer. Return only the summary.", f"Summarize in {lens.get(length,'1 paragraph')}:\n\n{text[:8000]}")
    return {"result": result}

@app.post("/ai-notes")
async def ai_notes(file: UploadFile = File(...), style: str = Form("bullet")):
    data = await read_file(file)
    text = extract_text_fitz(data, 25)
    if not text.strip(): raise HTTPException(400, "No readable text found.")
    style_map = {"bullet": "bullet point notes", "outline": "hierarchical outline", "cornell": "Cornell-style notes", "mindmap": "text-based mind map"}
    result = ai_call("You are an expert study notes creator. Return only formatted notes.", f"Create {style_map.get(style,'bullet points')}:\n\n{text[:9000]}")
    return {"result": result}

@app.post("/ai-quiz")
async def ai_quiz(file: UploadFile = File(...), count: int = Form(10), difficulty: str = Form("medium")):
    data = await read_file(file)
    text = extract_text_fitz(data, 20)
    if not text.strip(): raise HTTPException(400, "No readable text found.")
    count = max(3, min(count, 20))
    result = ai_call("You are a quiz creator. Return numbered questions with A–D options and mark correct answer with ✓.", f"Create {count} {difficulty} questions:\n\n{text[:8000]}")
    return {"result": result}

@app.post("/ai-keypoints")
async def ai_keypoints(file: UploadFile = File(...)):
    data = await read_file(file)
    text = extract_text_fitz(data, 20)
    if not text.strip(): raise HTTPException(400, "No readable text found.")
    result = ai_call("Extract the most important information as a numbered list.", f"Extract 8–12 key points:\n\n{text[:8000]}")
    return {"result": result}

@app.post("/ai-translate")
async def ai_translate(file: UploadFile = File(...), from_lang: str = Form("auto"), to_lang: str = Form("Urdu")):
    data = await read_file(file)
    text = extract_text_fitz(data, 15)
    if not text.strip(): raise HTTPException(400, "No readable text found.")
    src = f"from {from_lang}" if from_lang != "auto" else "(auto-detect)"
    result = ai_call(f"Translate accurately {src} to {to_lang}. Return only translated text.", f"Translate to {to_lang}:\n\n{text[:6000]}")
    return {"result": result}

@app.post("/ask-pdf")
async def ask_pdf(file: UploadFile = File(...), question: str = Form(...), history: str = Form("[]")):
    data = await read_file(file)
    text = extract_text_fitz(data, 30)
    if not text.strip(): raise HTTPException(400, "No readable text found.")
    try: hist = json.loads(history)[-6:]
    except: hist = []
    ctx = "\n".join(f"Q: {h['q']}\nA: {h['a']}" for h in hist if "q" in h and "a" in h)
    prompt = f"Document:\n{text[:9000]}\n\n"
    if ctx: prompt += f"Previous:\n{ctx}\n\n"
    prompt += f"Question: {question}"
    result = ai_call("Answer questions based strictly on the document. If not found, say so.", prompt, 1500)
    return {"result": result}


# ══════════════════════════════════════════
# PDF ORGANIZE
# ══════════════════════════════════════════

@app.post("/merge-pdf")
async def merge_pdf(files: List[UploadFile] = File(...)):
    if len(files) < 2: raise HTTPException(400, "Please upload at least 2 PDF files.")
    writer = fitz.open()
    for uf in files:
        data = await read_file(uf)
        try:
            doc = fitz.open(stream=data, filetype="pdf")
            writer.insert_pdf(doc); doc.close()
        except Exception as e:
            raise HTTPException(400, f"Could not read '{uf.filename}': {e}")
    out = io.BytesIO(); writer.save(out)
    return stream_file(out.getvalue(), "application/pdf", "merged.pdf")

@app.post("/split-pdf")
async def split_pdf(file: UploadFile = File(...), mode: str = Form("each"), start_page: int = Form(1), end_page: int = Form(1)):
    data = await read_file(file)
    doc = open_fitz(data); n = len(doc)
    if mode == "range":
        s = max(1, start_page) - 1; e = min(n, end_page) - 1
        if s > e: raise HTTPException(400, "Start page must be ≤ end page.")
        out_doc = fitz.open(); out_doc.insert_pdf(doc, from_page=s, to_page=e)
        buf = io.BytesIO(); out_doc.save(buf)
        return stream_file(buf.getvalue(), "application/pdf", f"{stem(file.filename)}_p{s+1}-{e+1}.pdf")
    else:
        zb = io.BytesIO()
        with zipfile.ZipFile(zb, "w", zipfile.ZIP_DEFLATED) as zf:
            for i in range(n):
                pg = fitz.open(); pg.insert_pdf(doc, from_page=i, to_page=i)
                buf = io.BytesIO(); pg.save(buf)
                zf.writestr(f"page_{i+1:03d}.pdf", buf.getvalue()); pg.close()
        return stream_file(zb.getvalue(), "application/zip", f"{stem(file.filename)}_pages.zip")

@app.post("/rotate-pdf")
async def rotate_pdf(file: UploadFile = File(...), angle: int = Form(90), pages: str = Form("all")):
    data = await read_file(file); doc = open_fitz(data); angle = angle % 360
    for i, page in enumerate(doc):
        apply = pages == "all" or (pages == "odd" and i % 2 == 0) or (pages == "even" and i % 2 == 1)
        if apply: page.set_rotation((page.rotation + angle) % 360)
    buf = io.BytesIO(); doc.save(buf, garbage=3, deflate=True)
    return stream_file(buf.getvalue(), "application/pdf", f"rotated_{stem(file.filename)}.pdf")

@app.post("/add-page-numbers")
async def add_page_numbers(file: UploadFile = File(...), position: str = Form("bottom-center"), format: str = Form("number"), start: int = Form(1)):
    data = await read_file(file); doc = open_fitz(data); n = len(doc)
    for i, page in enumerate(doc):
        num = i + start
        label = f"Page {num} of {n+start-1}" if format=="page-of" else (_to_roman(num) if format=="roman" else str(num))
        w, h = page.rect.width, page.rect.height; fs = 9
        tw = fitz.get_text_length(label, fontsize=fs)
        pos_map = {"bottom-center": ((w-tw)/2, h-18), "bottom-right": (w-tw-20, h-18), "top-center": ((w-tw)/2, 22), "top-right": (w-tw-20, 22)}
        x, y = pos_map.get(position, ((w-tw)/2, h-18))
        page.insert_text((x, y), label, fontsize=fs, color=(0.45, 0.45, 0.45))
    buf = io.BytesIO(); doc.save(buf, garbage=3, deflate=True)
    return stream_file(buf.getvalue(), "application/pdf", f"numbered_{stem(file.filename)}.pdf")


# ══════════════════════════════════════════
# COMPRESS
# ══════════════════════════════════════════

@app.post("/compress-pdf")
async def compress_pdf(file: UploadFile = File(...), level: str = Form("medium")):
    data = await read_file(file); orig = len(data); doc = open_fitz(data)
    buf = io.BytesIO()
    doc.save(buf, garbage=4, deflate=True, deflate_images=True, deflate_fonts=True, clean=True, linear=True)
    comp = buf.getvalue(); pct = round((1 - len(comp)/orig)*100, 1) if orig else 0
    return stream_file(comp, "application/pdf", f"compressed_{stem(file.filename)}.pdf",
                       {"X-Original-Size": str(orig), "X-Compressed-Size": str(len(comp)), "X-Savings-Pct": str(max(0,pct))})


# ══════════════════════════════════════════
# SECURITY
# ══════════════════════════════════════════

@app.post("/protect-pdf")
async def protect_pdf(file: UploadFile = File(...), password: str = Form(...)):
    data = await read_file(file)
    if not password: raise HTTPException(400, "Password is required.")
    doc = open_fitz(data)
    perm = fitz.PDF_PERM_PRINT | fitz.PDF_PERM_COPY
    buf = io.BytesIO()
    doc.save(buf, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw=password+"_owner", user_pw=password, permissions=perm)
    return stream_file(buf.getvalue(), "application/pdf", f"protected_{stem(file.filename)}.pdf")

@app.post("/unlock-pdf")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form("")):
    data = await read_file(file)
    doc = fitz.open(stream=data, filetype="pdf")
    if doc.is_encrypted:
        if not doc.authenticate(password): raise HTTPException(400, "Wrong password.")
    tmp = "/tmp/_rundocs_unlock.pdf"
    doc.save(tmp, encryption=fitz.PDF_ENCRYPT_NONE)
    with open(tmp, "rb") as f: result = f.read()
    os.remove(tmp)
    return stream_file(result, "application/pdf", f"unlocked_{stem(file.filename)}.pdf")


# ══════════════════════════════════════════
# WATERMARK FIX — uses shapes layer properly
# ══════════════════════════════════════════

@app.post("/add-watermark")
async def add_watermark(
    file: UploadFile = File(...),
    text: str = Form("CONFIDENTIAL"),
    opacity: float = Form(0.2),
    position: str = Form("center"),
):
    data = await read_file(file)
    doc = open_fitz(data)
    opacity = max(0.05, min(opacity, 0.9))

    for page in doc:
        w, h = page.rect.width, page.rect.height

        if position == "center":
            # Diagonal center watermark
            font_size = min(w, h) * 0.07
            # Create a transparent overlay using insert_text with proper color+opacity
            tw = fitz.get_text_length(text, fontname="helv", fontsize=font_size)
            # Center point
            cx, cy = w / 2, h / 2
            # Use insert_text with rotation for diagonal effect
            page.insert_text(
                fitz.Point(cx - tw/2, cy),
                text,
                fontname="helv",
                fontsize=font_size,
                color=(0.6, 0.1, 0.1),
                rotate=45,
                overlay=True,
            )
            # Apply opacity via transparency by drawing a white rect on top won't work
            # Instead use the correct approach: draw text with alpha using Shape
            # Re-do using shape for proper opacity support
            page_rect = page.rect
            shape = page.new_shape()
            # Calculate rotated text position for center diagonal
            shape.insert_text(
                fitz.Point(cx - tw * 0.6, cy + tw * 0.35),
                text,
                fontname="helv",
                fontsize=font_size * 0.85,
                color=(0.55, 0.08, 0.08),
                rotate=45,
            )
            shape.finish(
                fill=None,
                color=(0.55, 0.08, 0.08),
                fill_opacity=opacity,
                stroke_opacity=opacity,
            )
            shape.commit(overlay=True)

        elif position == "bottom-right":
            font_size = min(w, h) * 0.042
            tw = fitz.get_text_length(text, fontname="helv", fontsize=font_size)
            shape = page.new_shape()
            shape.insert_text(
                fitz.Point(w - tw - 22, h - 22),
                text,
                fontname="helv",
                fontsize=font_size,
                color=(0.35, 0.35, 0.35),
            )
            shape.finish(color=(0.35,0.35,0.35), fill_opacity=opacity, stroke_opacity=opacity)
            shape.commit(overlay=True)

        else:  # bottom-center
            font_size = min(w, h) * 0.042
            tw = fitz.get_text_length(text, fontname="helv", fontsize=font_size)
            shape = page.new_shape()
            shape.insert_text(
                fitz.Point((w - tw) / 2, h - 22),
                text,
                fontname="helv",
                fontsize=font_size,
                color=(0.35, 0.35, 0.35),
            )
            shape.finish(color=(0.35,0.35,0.35), fill_opacity=opacity, stroke_opacity=opacity)
            shape.commit(overlay=True)

    buf = io.BytesIO()
    doc.save(buf, garbage=3, deflate=True)
    return stream_file(buf.getvalue(), "application/pdf", f"watermarked_{stem(file.filename)}.pdf")


# ══════════════════════════════════════════
# CONVERT
# ══════════════════════════════════════════

@app.post("/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    data = await read_file(file)
    if not _DOCX_OK: raise HTTPException(500, "Word library not available.")
    doc_out = Document()
    doc_out.styles["Normal"].font.name = "Calibri"
    doc_out.styles["Normal"].font.size = Pt(11)
    fitz_doc = open_fitz(data)
    for page_num in range(len(fitz_doc)):
        page = fitz_doc[page_num]
        if page_num > 0: doc_out.add_page_break()
        for block in page.get_text("dict")["blocks"]:
            if block.get("type") != 0: continue
            for line in block.get("lines", []):
                line_text = " ".join(s.get("text","") for s in line.get("spans",[])).strip()
                if not line_text: continue
                sizes = [s.get("size",11) for s in line.get("spans",[])]
                avg = sum(sizes)/len(sizes) if sizes else 11
                if avg > 16: doc_out.add_heading(line_text, level=1)
                elif avg > 13: doc_out.add_heading(line_text, level=2)
                else: doc_out.add_paragraph(line_text)
    buf = io.BytesIO(); doc_out.save(buf)
    return stream_file(buf.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document", f"{stem(file.filename)}.docx")


# ══════════════════════════════════════════════════════════════════
# PDF TO EXCEL — SMART GROUPING FIX (header-based, NOT page-based)
# ══════════════════════════════════════════════════════════════════

@app.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...), mode: str = Form("smart")):
    data = await read_file(file)
    if not _XL_OK:      raise HTTPException(500, "Excel library not available.")
    if not _PLUMBER_OK: raise HTTPException(500, "PDF processing library not available.")

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    # ── Styles ────────────────────────────────────────────────────────────────
    hdr_fill  = PatternFill("solid", fgColor="4F6EF7")
    hdr_font  = Font(bold=True, color="FFFFFF", size=11, name="Calibri")
    hdr_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    alt_fill  = PatternFill("solid", fgColor="EEF1FF")
    thin      = Side(style="thin", color="C5CEFF")
    border    = Border(left=thin, right=thin, top=thin, bottom=thin)

    # ── Helpers ───────────────────────────────────────────────────────────────
    def clean(v):
        return re.sub(r"\s+", " ", str(v or "").strip())

    def norm_key(row):
        """Canonical header key — lowercase, collapsed spaces, no punctuation."""
        return tuple(re.sub(r"[^\w]", "", clean(c).lower()) for c in row)

    def fuzzy_match(k1, k2):
        """Two header keys match if ≥75% of tokens are identical (handles minor OCR drift)."""
        if k1 == k2:
            return True
        if not k1 or not k2:
            return False
        shorter, longer = (k1, k2) if len(k1) <= len(k2) else (k2, k1)
        if len(shorter) == 0:
            return False
        hits = sum(1 for t in shorter if t in longer)
        return hits / len(shorter) >= 0.75

    def looks_like_header(row):
        """True when row is mostly text labels, not numbers."""
        cells = [clean(c) for c in row if clean(c)]
        if len(cells) < 2:
            return False
        numeric = sum(1 for c in cells if re.match(r"^[\d\s\.\,\-\$\%\+\(\)\/\#]+$", c))
        return (numeric / len(cells)) < 0.4

    def find_header_idx(tbl):
        for i, row in enumerate(tbl[:5]):
            if looks_like_header(row):
                return i
        return 0

    def make_sheet_name(existing):
        n = len(existing) + 1
        return f"Table {n}"

    def coerce(val):
        s = re.sub(r"[,\$%]", "", clean(val))
        if not s:
            return ""
        try:
            return int(s) if "." not in s else round(float(s), 6)
        except (ValueError, TypeError):
            return val

    def find_matching_group(hkey, table_groups):
        """Find existing group whose key fuzzy-matches hkey."""
        for k in table_groups:
            if fuzzy_match(k, hkey):
                return k
        return None

    # ── Extraction ────────────────────────────────────────────────────────────
    table_groups: dict = {}
    text_lines:   list = []

    TABLE_SETTINGS = [
        # Strategy 1: explicit borders
        {"vertical_strategy": "lines_strict", "horizontal_strategy": "lines_strict",
         "snap_tolerance": 4, "join_tolerance": 4, "edge_min_length": 10,
         "min_words_vertical": 1, "min_words_horizontal": 1},
        # Strategy 2: text-aligned (borderless)
        {"vertical_strategy": "text", "horizontal_strategy": "text",
         "text_tolerance": 4, "intersection_tolerance": 4},
        # Strategy 3: looser lines
        {"vertical_strategy": "lines", "horizontal_strategy": "lines",
         "snap_tolerance": 6, "join_tolerance": 6, "edge_min_length": 6},
    ]

    def extract_page_tables(page):
        """Try each strategy, return first that yields valid tables."""
        for s in TABLE_SETTINGS:
            try:
                tbls = page.extract_tables(s) or []
                tbls = [[[clean(c) for c in row] for row in t]
                        for t in tbls
                        if t and len(t) >= 2
                        and max((len(r) for r in t), default=0) >= 2]
                if tbls:
                    return tbls
            except Exception:
                continue
        return []

    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            # Pass 1: collect ALL tables from ALL pages with their page number
            all_tables = []  # list of (page_num, table)
            for pn, page in enumerate(pdf.pages, start=1):
                if mode in ("smart", "tables"):
                    for tbl in extract_page_tables(page):
                        all_tables.append((pn, tbl))
                if mode == "text":
                    raw = page.extract_text() or ""
                    for ln in raw.split("\n"):
                        if ln.strip():
                            text_lines.append((pn, ln.strip()))

            # Pass 2: group tables by header across all pages
            # table_groups key → {sheet_name, header, rows, col_count}
            for pn, tbl in all_tables:
                if not tbl:
                    continue

                hdr_idx = find_header_idx(tbl)
                raw_hdr = tbl[hdr_idx]

                if not any(raw_hdr):
                    continue

                hkey     = norm_key(raw_hdr)
                disp_hdr = raw_hdr
                col_count = len(raw_hdr)

                # Collect data rows — skip empty, skip repeated header rows
                data_rows = []
                for row in tbl[hdr_idx + 1:]:
                    if not any(v for v in row):
                        continue
                    rk = norm_key(row)
                    # Skip if this row is a repeated header
                    if fuzzy_match(rk, hkey):
                        continue
                    data_rows.append(row)

                if not data_rows:
                    continue

                # Find if a matching group already exists
                matched = find_matching_group(hkey, table_groups)

                if matched:
                    # Merge into existing sheet
                    grp = table_groups[matched]
                    # Pad/trim rows to match existing col count
                    c = grp["col_count"]
                    for row in data_rows:
                        grp["rows"].append((row + [""] * c)[:c])
                else:
                    # New group → new sheet
                    sname = make_sheet_name(
                        [g["sheet_name"] for g in table_groups.values()]
                    )
                    table_groups[hkey] = {
                        "sheet_name": sname,
                        "header":     disp_hdr,
                        "rows":       data_rows,
                        "col_count":  col_count,
                    }

            # If no tables found at all, fallback to text extraction
            if not table_groups and not text_lines and mode != "text":
                with pdfplumber.open(io.BytesIO(data)) as pdf2:
                    for pn, page in enumerate(pdf2.pages, start=1):
                        raw = page.extract_text() or ""
                        for ln in raw.split("\n"):
                            if ln.strip():
                                text_lines.append((pn, ln.strip()))

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"PDF processing failed: {e}")

    # ── Write workbook ────────────────────────────────────────────────────────
    if table_groups:
        for grp in table_groups.values():
            ws = wb.create_sheet(title=grp["sheet_name"])
            ws.freeze_panes = "A2"

            cols = grp["col_count"]
            # Header row
            for ci, h in enumerate(grp["header"], start=1):
                cell = ws.cell(row=1, column=ci, value=h)
                cell.fill = hdr_fill; cell.font = hdr_font
                cell.alignment = hdr_align; cell.border = border

            # Data rows
            for ri, row in enumerate(grp["rows"], start=2):
                padded = (row + [""] * cols)[:cols]
                for ci, val in enumerate(padded, start=1):
                    cell = ws.cell(row=ri, column=ci, value=coerce(val))
                    cell.border = border
                    if ri % 2 == 0:
                        cell.fill = alt_fill

            # Auto column width
            for col in ws.columns:
                mx = max((len(str(c.value or "")) for c in col), default=10)
                ws.column_dimensions[col[0].column_letter].width = min(max(mx + 3, 12), 60)

    elif text_lines:
        ws = wb.create_sheet(title="Extracted Text")
        ws.freeze_panes = "A2"
        for ci, h in enumerate(["Page", "Text"], start=1):
            cell = ws.cell(row=1, column=ci, value=h)
            cell.fill = hdr_fill; cell.font = hdr_font
            cell.alignment = hdr_align; cell.border = border
        for ri, (pg, ln) in enumerate(text_lines, start=2):
            ws.cell(row=ri, column=1, value=pg).border = border
            ws.cell(row=ri, column=2, value=ln).border = border
            if ri % 2 == 0:
                ws.cell(ri, 1).fill = alt_fill
                ws.cell(ri, 2).fill = alt_fill
        ws.column_dimensions["A"].width = 8
        ws.column_dimensions["B"].width = 90

    else:
        ws = wb.create_sheet(title="No Data Found")
        ws.cell(row=1, column=1, value="No tables or text could be extracted from this PDF.")
        ws.cell(row=2, column=1, value="Tip: PDF must contain selectable text (not scanned images).")

    out = io.BytesIO(); wb.save(out)
    return stream_file(out.getvalue(),
                       "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                       f"{stem(file.filename)}.xlsx")


@app.post("/pdf-to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...), dpi: int = Form(150)):
    data = await read_file(file); doc = open_fitz(data)
    dpi = max(72, min(dpi, 300)); mat = fitz.Matrix(dpi/72, dpi/72)
    zb = io.BytesIO()
    with zipfile.ZipFile(zb, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            zf.writestr(f"page_{i+1:03d}.jpg", pix.tobytes("jpeg"))
    return stream_file(zb.getvalue(), "application/zip", f"{stem(file.filename)}_images.zip")


@app.post("/jpg-to-pdf")
async def jpg_to_pdf(files: List[UploadFile] = File(...)):
    doc = fitz.open()
    for uf in files:
        data = await read_file(uf)
        try:
            ext = (uf.content_type or "jpeg").split("/")[-1]
            img_doc = fitz.open(stream=data, filetype=ext)
            page = doc.new_page(width=img_doc[0].rect.width, height=img_doc[0].rect.height)
            page.show_pdf_page(page.rect, img_doc, 0); img_doc.close()
        except Exception:
            try:
                page = doc.new_page(width=595, height=842)
                page.insert_image(fitz.Rect(0,0,595,842), stream=data)
            except Exception as e:
                raise HTTPException(400, f"Could not process '{uf.filename}': {e}")
    buf = io.BytesIO(); doc.save(buf)
    return stream_file(buf.getvalue(), "application/pdf", "images.pdf")


@app.post("/pdf-to-pptx")
async def pdf_to_pptx(file: UploadFile = File(...)):
    data = await read_file(file)
    if not _PPTX_OK: raise HTTPException(500, "PowerPoint library not available.")
    doc = open_fitz(data); prs = Presentation()
    prs.slide_width = PInches(10); prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(1.5,1.5))
        img = pix.tobytes("png")
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(io.BytesIO(img), PInches(0), PInches(0), width=prs.slide_width, height=prs.slide_height)
    buf = io.BytesIO(); prs.save(buf)
    return stream_file(buf.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation", f"{stem(file.filename)}.pptx")


@app.post("/pdf-to-text")
async def pdf_to_text(file: UploadFile = File(...)):
    data = await read_file(file)
    if _PLUMBER_OK:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            lines = []
            for i, page in enumerate(pdf.pages):
                lines.append(f"─── Page {i+1} ───\n")
                lines.append(page.extract_text() or "[No text]")
                lines.append("\n")
            text = "\n".join(lines)
    else:
        text = extract_text_fitz(data)
    return stream_file(text.encode("utf-8"), "text/plain", f"{stem(file.filename)}.txt")


@app.post("/ocr-check")
async def ocr_check(file: UploadFile = File(...)):
    data = await read_file(file); doc = open_fitz(data)
    pages_checked = min(len(doc), 3)
    total_chars = sum(len(doc[i].get_text()) for i in range(pages_checked))
    return {"is_scanned": total_chars < 80, "chars_found": total_chars, "pages_checked": pages_checked}


# ══════════════════════════════════════════════════════════════════
# MERGED FROM additional_endpoints.py — powers word-to-pdf.html,
# excel-to-pdf.html, delete-pages.html, reorder-pages.html, sign-pdf.html
# ══════════════════════════════════════════════════════════════════

def _libreoffice_available() -> bool:
    return shutil.which("libreoffice") is not None or shutil.which("soffice") is not None


def _convert_with_libreoffice(input_bytes: bytes, input_ext: str, out_ext: str = "pdf") -> bytes:
    """Convert a file to PDF (or other format) using headless LibreOffice."""
    binary = shutil.which("soffice") or shutil.which("libreoffice")
    if not binary:
        raise HTTPException(500, "Conversion engine not available on server.")

    with tempfile.TemporaryDirectory() as tmp:
        in_path = os.path.join(tmp, f"input.{input_ext}")
        with open(in_path, "wb") as f:
            f.write(input_bytes)

        cmd = [
            binary, "--headless", "--norestore", "--invisible",
            "--convert-to", out_ext, "--outdir", tmp, in_path,
        ]
        try:
            subprocess.run(cmd, check=True, timeout=60, capture_output=True)
        except subprocess.TimeoutExpired:
            raise HTTPException(504, "Conversion timed out. Try a smaller file.")
        except subprocess.CalledProcessError as e:
            raise HTTPException(500, f"Conversion failed: {e.stderr.decode(errors='ignore')[:200]}")

        out_path = os.path.join(tmp, f"input.{out_ext}")
        if not os.path.exists(out_path):
            raise HTTPException(500, "Conversion produced no output file.")
        with open(out_path, "rb") as f:
            return f.read()


# ══════════════════════════════════════════
# WORD TO PDF
# ══════════════════════════════════════════

@app.post("/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    data = await read_file(file)
    ext = "docx" if file.filename.lower().endswith(".docx") else "doc"

    if _libreoffice_available():
        pdf_bytes = _convert_with_libreoffice(data, ext, "pdf")
        return stream_file(pdf_bytes, "application/pdf", f"{stem(file.filename)}.pdf")

    # Fallback: extract text with python-docx and render a simple PDF with fitz
    if not _DOCX_OK or not _FITZ_OK:
        raise HTTPException(500, "Conversion engine not available. Contact support.")
    try:
        word_doc = Document(io.BytesIO(data))
    except Exception as e:
        raise HTTPException(400, f"Could not read Word file: {e}")

    pdf_doc = fitz.open()
    page = pdf_doc.new_page()
    y = 50
    for para in word_doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if y > 780:
            page = pdf_doc.new_page()
            y = 50
        page.insert_text((50, y), text, fontsize=11)
        y += 16
    buf = io.BytesIO()
    pdf_doc.save(buf)
    return stream_file(buf.getvalue(), "application/pdf", f"{stem(file.filename)}.pdf")


# ══════════════════════════════════════════
# EXCEL TO PDF
# ══════════════════════════════════════════

@app.post("/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    data = await read_file(file)
    ext = "xlsx" if file.filename.lower().endswith(".xlsx") else "xls"

    if _libreoffice_available():
        pdf_bytes = _convert_with_libreoffice(data, ext, "pdf")
        return stream_file(pdf_bytes, "application/pdf", f"{stem(file.filename)}.pdf")

    # Fallback: render Excel cell values as a simple table in a PDF
    if not _XL_OK or not _FITZ_OK:
        raise HTTPException(500, "Conversion engine not available. Contact support.")
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    except Exception as e:
        raise HTTPException(400, f"Could not read Excel file: {e}")

    pdf_doc = fitz.open()
    for ws in wb.worksheets:
        page = pdf_doc.new_page()
        y = 50
        page.insert_text((40, 35), f"Sheet: {ws.title}", fontsize=13)
        for row in ws.iter_rows(values_only=True):
            if y > 780:
                page = pdf_doc.new_page()
                y = 50
            line = "  |  ".join(str(c) if c is not None else "" for c in row)
            page.insert_text((40, y), line[:140], fontsize=9)
            y += 14
    buf = io.BytesIO()
    pdf_doc.save(buf)
    return stream_file(buf.getvalue(), "application/pdf", f"{stem(file.filename)}.pdf")


# ══════════════════════════════════════════
# DELETE PAGES  (e.g. "1,3,5" or "2-4" or "1,3-5,8")
# ══════════════════════════════════════════

def _parse_page_spec(spec: str, total_pages: int) -> set:
    """Parses '1,3,5' or '2-4' or '1,3-5,8' into a 0-indexed set of page numbers."""
    pages = set()
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            try:
                a, b = part.split("-")
                a, b = int(a.strip()), int(b.strip())
                for p in range(min(a, b), max(a, b) + 1):
                    if 1 <= p <= total_pages:
                        pages.add(p - 1)
            except ValueError:
                continue
        else:
            try:
                p = int(part)
                if 1 <= p <= total_pages:
                    pages.add(p - 1)
            except ValueError:
                continue
    return pages


@app.post("/delete-pages")
async def delete_pages(file: UploadFile = File(...), pages: str = Form(...)):
    data = await read_file(file)
    doc = open_fitz(data)
    total = len(doc)

    to_delete = _parse_page_spec(pages, total)
    if not to_delete:
        raise HTTPException(400, "No valid page numbers found. Use format like 1,3,5 or 2-4.")
    if len(to_delete) >= total:
        raise HTTPException(400, "Cannot delete all pages from the document.")

    doc.delete_pages(sorted(to_delete))
    buf = io.BytesIO()
    doc.save(buf, garbage=3, deflate=True)
    return stream_file(buf.getvalue(), "application/pdf", f"edited_{stem(file.filename)}.pdf")


# ══════════════════════════════════════════
# REORDER PAGES  (e.g. "3,1,2,4")
# ══════════════════════════════════════════

@app.post("/reorder-pages")
async def reorder_pages(file: UploadFile = File(...), order: str = Form(...)):
    data = await read_file(file)
    doc = open_fitz(data)
    total = len(doc)

    try:
        new_order = [int(x.strip()) - 1 for x in order.split(",") if x.strip()]
    except ValueError:
        raise HTTPException(400, "Invalid page order format. Use e.g. 3,1,2,4")

    if sorted(new_order) != list(range(total)):
        raise HTTPException(
            400,
            f"Page order must include every page exactly once (1 to {total}).",
        )

    new_doc = fitz.open()
    for idx in new_order:
        new_doc.insert_pdf(doc, from_page=idx, to_page=idx)

    buf = io.BytesIO()
    new_doc.save(buf, garbage=3, deflate=True)
    return stream_file(buf.getvalue(), "application/pdf", f"reordered_{stem(file.filename)}.pdf")


# ══════════════════════════════════════════
# SIGN PDF  (text-based signature stamp)
# ══════════════════════════════════════════

@app.post("/sign-pdf")
async def sign_pdf(
    file: UploadFile = File(...),
    text: str = Form(...),
    position: str = Form("bottom-right"),
    pages: str = Form("last"),
):
    data = await read_file(file)
    doc = open_fitz(data)
    total = len(doc)

    if pages == "all":
        target_pages = range(total)
    elif pages == "first":
        target_pages = [0]
    else:  # "last"
        target_pages = [total - 1]

    for i in target_pages:
        page = doc[i]
        w, h = page.rect.width, page.rect.height
        font_size = 14
        tw = fitz.get_text_length(text, fontname="helv", fontsize=font_size)

        pos_map = {
            "bottom-right": (w - tw - 40, h - 45),
            "bottom-center": ((w - tw) / 2, h - 45),
            "bottom-left": (40, h - 45),
        }
        x, y = pos_map.get(position, (w - tw - 40, h - 45))

        # Draw a subtle line above the signature for a more "signed" look
        page.draw_line(
            fitz.Point(x, y - 6), fitz.Point(x + tw, y - 6),
            color=(0.2, 0.2, 0.2), width=0.6,
        )
        page.insert_text(
            (x, y), text,
            fontname="helv", fontsize=font_size,
            color=(0.1, 0.1, 0.5),
        )
        page.insert_text(
            (x, y + 14), "Digitally signed via RunDocs",
            fontname="helv", fontsize=7,
            color=(0.5, 0.5, 0.5),
        )

    buf = io.BytesIO()
    doc.save(buf, garbage=3, deflate=True)
    return stream_file(buf.getvalue(), "application/pdf", f"signed_{stem(file.filename)}.pdf")

# ══════════════════════════════════════════════════════════════════
# HTML → PDF
# ══════════════════════════════════════════════════════════════════
@app.post("/html-to-pdf")
async def html_to_pdf(file: UploadFile = File(...)):
    data = await read_file(file)
    fname = (file.filename or "page.html").lower()

    # Try WeasyPrint first (best quality, handles CSS/images)
    try:
        from weasyprint import HTML as WH, CSS
        html_str = data.decode("utf-8", errors="replace")
        pdf_bytes = WH(string=html_str, base_url=None).write_pdf()
        return stream_file(pdf_bytes, "application/pdf", f"{stem(file.filename)}.pdf")
    except ImportError:
        pass
    except Exception:
        pass

    # Fallback: LibreOffice
    if _libreoffice_available():
        try:
            pdf_bytes = _convert_with_libreoffice(data, "html", "pdf")
            return stream_file(pdf_bytes, "application/pdf", f"{stem(file.filename)}.pdf")
        except Exception:
            pass

    # Last resort: embed HTML text into a simple PDF via fitz
    try:
        html_str = data.decode("utf-8", errors="replace")
        import re as _re
        text = _re.sub(r"<[^>]+>", " ", html_str)
        text = _re.sub(r"\s+", " ", text).strip()
        doc = fitz.open()
        page = doc.new_page()
        rect = fitz.Rect(50, 50, 562, 742)
        page.insert_textbox(rect, text, fontsize=11, fontname="helv")
        buf = io.BytesIO()
        doc.save(buf)
        return stream_file(buf.getvalue(), "application/pdf", f"{stem(file.filename)}.pdf")
    except Exception as e:
        raise HTTPException(500, f"HTML to PDF conversion failed: {e}")


# ══════════════════════════════════════════════════════════════════
# PowerPoint → PDF  (pptx / ppt)
# ══════════════════════════════════════════════════════════════════
@app.post("/pptx-to-pdf")
async def pptx_to_pdf(file: UploadFile = File(...)):
    data = await read_file(file)
    fname = (file.filename or "presentation.pptx").lower()
    ext   = "pptx" if fname.endswith(".pptx") else "ppt"

    # LibreOffice (best quality, preserves images/fonts)
    if _libreoffice_available():
        try:
            pdf_bytes = _convert_with_libreoffice(data, ext, "pdf")
            return stream_file(pdf_bytes, "application/pdf", f"{stem(file.filename)}.pdf")
        except Exception:
            pass

    # Fallback: extract text from each slide via python-pptx and write to PDF
    try:
        from pptx import Presentation as _Prs
        prs = _Prs(io.BytesIO(data))
    except Exception as e:
        raise HTTPException(400, f"Could not read PowerPoint file: {e}")

    try:
        doc = fitz.open()
        for slide_num, slide in enumerate(prs.slides, start=1):
            page = doc.new_page(width=792, height=612)   # landscape A4
            lines = [f"Slide {slide_num}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
            text = "\n".join(lines)
            rect = fitz.Rect(40, 40, 752, 572)
            page.insert_textbox(rect, text, fontsize=13, fontname="helv")
        buf = io.BytesIO()
        doc.save(buf, garbage=3, deflate=True)
        return stream_file(buf.getvalue(), "application/pdf", f"{stem(file.filename)}.pdf")
    except Exception as e:
        raise HTTPException(500, f"PowerPoint to PDF conversion failed: {e}")
        # ── Paddle Webhook ──────────────────────────────────────────
import os as _os
import json as _json
from supabase import create_client as _create_client

@app.post("/paddle/webhook")
async def paddle_webhook(request: Request):
    try:
        body = await request.body()
        payload = _json.loads(body)
        event_type = payload.get("event_type", "")
        data = payload.get("data", {})

        sb = _create_client(
            _os.getenv("SUPABASE_URL"),
            _os.getenv("SUPABASE_SERVICE_KEY")
        )

        if event_type in ["subscription.created", "subscription.updated"]:
            custom_data = data.get("custom_data", {})
            user_id = custom_data.get("user_id")
            plan = custom_data.get("plan", "pro")
            paddle_subscription_id = data.get("id")
            paddle_customer_id = data.get("customer_id")
            status = data.get("status", "active")
            billing = data.get("billing_cycle", {}).get("interval", "month")
            billing_period = "yearly" if "year" in str(billing) else "monthly"
            if user_id:
                sb.table("Subscriptions").upsert({
                    "User_id": user_id,
                    "paddle_subscription_id": paddle_subscription_id,
                    "paddle_customer_id": paddle_customer_id,
                    "Plan": plan,
                    "Status": status,
                    "billing_period": billing_period,
                }, on_conflict="paddle_subscription_id").execute()

        elif event_type == "subscription.canceled":
            paddle_subscription_id = data.get("id")
            if paddle_subscription_id:
                sb.table("Subscriptions").update({
                    "Status": "canceled",
                    "Plan": "free"
                }).eq("paddle_subscription_id", paddle_subscription_id).execute()

        return {"status": "ok"}
    except Exception as e:
        return {"status": "error", "detail": str(e)}

